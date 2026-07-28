#!/usr/bin/env python3
"""
VigyanLLM Razorpay Payment Integration — PostgreSQL Version
==============================================================
Hybrid Subscription + Token Pack model with full cost tracking.

Security hardening:
  PAY-01 FIX: verify_payment() now calls Razorpay API GET /payments/{id}
              after HMAC verification to confirm payment is actually captured
              server-side. Prevents forged HMAC claims.
  PAY-04 FIX: RAZORPAY_WEBHOOK_SECRET must be set explicitly in env;
              it NO LONGER falls back to RAZORPAY_KEY_SECRET.
  BUG-38 FIX: Removed hardcoded test phone number from Razorpay orders.

Endpoints:
  POST /api/payments/create-order    — Create Razorpay order
  POST /api/payments/verify-payment  — Verify signature & confirm via API
  POST /api/payments/webhook         — Razorpay webhook (server-to-server)
  GET  /api/payments/pricing         — Public pricing data
  GET  /api/payments/token-balance   — User's token balance & subscription
  GET  /api/payments/financial-summary — Admin: P&L and ROI views
"""

import hashlib
import hmac
import json
import logging
import os
import time

import razorpay
from flask import Blueprint, Response, g, jsonify, request

from .database import (
    execute,
    fetch_all,
    fetch_one,
    get_db,
    get_db_standalone,
    put_db_standalone,
)
from .pg_auth import check_usage, log_action, require_admin, require_auth
from .price_registry import (
    FREE_TRIAL_RUNS,
    PRICE_REGISTRY,
    TOPUP_PRICE_INR,
    ACADEMIC_DISCOUNT_PCT,
    PLAN_REGISTRY,
    get_academic_price,
    get_amount_paise,
    get_designs_for_product,
    get_dock_runs_for_product,
    validate_order_request,
)
from .security import validate_quantity

logger = logging.getLogger("primerforge.payments")

payment_bp = Blueprint("payments", __name__)

# ── Razorpay Configuration ────────────────────────────────────────────────
RAZORPAY_KEY_ID = os.environ.get("RAZORPAY_KEY_ID", "")
RAZORPAY_KEY_SECRET = os.environ.get("RAZORPAY_KEY_SECRET", "")
RAZORPAY_WEBHOOK_SECRET = os.environ.get("RAZORPAY_WEBHOOK_SECRET", "")

if not RAZORPAY_KEY_ID or not RAZORPAY_KEY_SECRET:
    logger.warning("RAZORPAY_KEY_ID or RAZORPAY_KEY_SECRET not set — payment endpoints will fail")

# PAY-04 FIX: RAZORPAY_WEBHOOK_SECRET must be set independently.
# It MUST NOT fall back to RAZORPAY_KEY_SECRET. If both were the same value,
# a compromised key secret would also compromise webhook validation.
# The webhook secret is configured in Razorpay Dashboard → Webhooks → Secret.
if not RAZORPAY_WEBHOOK_SECRET:
    logger.error(
        "RAZORPAY_WEBHOOK_SECRET is not set. Razorpay webhooks will be REJECTED. "
        "Set this in your .env file (must differ from RAZORPAY_KEY_SECRET). "
        "Generate with: python3 -c \"import secrets; print(secrets.token_hex(32))\""
    )
    # Do NOT fall back to key secret — fail loudly so operators notice

rz_client = razorpay.Client(auth=(RAZORPAY_KEY_ID, RAZORPAY_KEY_SECRET)) if RAZORPAY_KEY_ID else None


# ── Helper: Verify Razorpay Signature ─────────────────────────────────────

def _verify_signature(order_id: str, payment_id: str, signature: str) -> bool:
    """
    Verify the Razorpay payment signature using HMAC-SHA256.

    Signature format: HMAC-SHA256(order_id + '|' + payment_id, key_secret)
    Uses hmac.compare_digest() for constant-time comparison to prevent
    timing oracle attacks.

    Args:
        order_id:   Razorpay order ID (razorpay_order_id from client).
        payment_id: Razorpay payment ID (razorpay_payment_id from client).
        signature:  Hex HMAC signature (razorpay_signature from client).

    Returns:
        True if signature matches, False otherwise.
    """
    message = f"{order_id}|{payment_id}"
    expected = hmac.new(
        RAZORPAY_KEY_SECRET.encode("utf-8"),
        message.encode("utf-8"),
        hashlib.sha256
    ).hexdigest()
    return hmac.compare_digest(expected, signature)


def _confirm_payment_server_side(payment_id: str) -> bool:
    """
    Confirm payment status via Razorpay API (PAY-01 FIX).

    HMAC signature verification alone is insufficient: a sophisticated attacker
    who knows the HMAC secret could forge a valid signature for a payment that
    was never actually captured. This function calls the Razorpay API directly
    to verify the payment is in 'captured' status before crediting tokens.

    Args:
        payment_id: Razorpay payment ID to verify.

    Returns:
        True if payment.status == 'captured', False otherwise.
        Returns False on API error (fail-closed).
    """
    if not rz_client:
        logger.error("Razorpay client not initialised — cannot confirm payment server-side")
        return False
    try:
        payment = rz_client.payment.fetch(payment_id)
        status = payment.get("status", "")
        if status != "captured":
            logger.warning(
                "PAY-01: Razorpay API returned status '%s' for payment %s (expected 'captured')",
                status, payment_id,
            )
            return False
        return True
    except Exception as e:
        logger.error("Razorpay API confirmation failed for payment %s: %s", payment_id, e)
        # Fail-closed: do not credit tokens if we cannot confirm with Razorpay
        return False


def _credit_tokens_atomic(user_id: int, order_id: str, product_id: str,
                           quantity: int = 1, payment_id: str = "") -> int:
    """
    Atomically credit designs after payment capture. Idempotent.
    - Subscriptions: activate plan + set monthly quota
    - Top-ups: add to balance
    Returns designs credited (0 if already processed).
    """
    db = get_db()
    cur = db.cursor()

    # Conditional update: only if status is still 'initiated' or 'authorized'
    cur.execute(
        """UPDATE payments SET status = 'captured', captured_at = NOW(),
                gateway_payment_id = COALESCE(NULLIF(%s, ''), gateway_payment_id)
           WHERE gateway_order_id = %s AND status IN ('initiated', 'authorized')
           RETURNING id, product_type, tokens_purchased""",
        (payment_id, order_id)
    )
    row = cur.fetchone()

    if not row:
        cur.close()
        return 0  # Already processed (idempotent)

    designs = get_designs_for_product(product_id, quantity)
    dock_runs = get_dock_runs_for_product(product_id, quantity)

    if product_id in PRICE_REGISTRY:
        # Subscription plan — activate and set monthly quotas
        product = PRICE_REGISTRY[product_id]
        expiry_interval = "24 hours" if product.period == "daily" else "30 days"
        cur.execute(
            f"""INSERT INTO subscriptions (user_id, is_active, plan_id, plan_type, monthly_quota,
                   dock_monthly_quota, quota_used, dock_quota_used,
                   started_at, expires_at, last_renewed_at, max_seats, quota_reset_at)
               VALUES (%s, TRUE, %s, %s, %s, %s, 0, 0,
                       NOW(), NOW() + INTERVAL '{expiry_interval}', NOW(), %s, NOW() + INTERVAL '{expiry_interval}')
               ON CONFLICT (user_id) DO UPDATE SET
                 is_active = TRUE,
                 plan_id = %s,
                 plan_type = %s,
                 monthly_quota = %s,
                 dock_monthly_quota = %s,
                 quota_used = 0,
                 dock_quota_used = 0,
                 expires_at = NOW() + INTERVAL '{expiry_interval}',
                 last_renewed_at = NOW(),
                 max_seats = %s,
                 quota_reset_at = NOW() + INTERVAL '{expiry_interval}'""",
            (user_id, product_id, product_id, designs, dock_runs, product.max_seats,
             product_id, product_id, designs, dock_runs, product.max_seats)
        )
    else:
        # Top-up — add to balance directly
        cur.execute(
            """UPDATE token_balances
               SET balance = balance + %s,
                   total_purchased = total_purchased + %s,
                   last_credited_at = NOW()
               WHERE user_id = %s""",
            (designs, designs, user_id)
        )

    db.commit()
    cur.close()
    return designs


# ══════════════════════════════════════════════════════════════════════════
# ENDPOINTS
# ══════════════════════════════════════════════════════════════════════════

@payment_bp.route("/api/payments/pricing", methods=["GET"])
def get_pricing():
    """Public endpoint: return all plan pricing from registry."""
    plans = []
    for cfg in PLAN_REGISTRY.values():
        if not cfg.is_active:
            continue
        plans.append({
            "plan_id": cfg.plan_id,
            "display_name": cfg.display_name,
            "tier": cfg.tier.value,
            "billing": cfg.billing.value,
            "price_inr": cfg.price_inr,
            "academic_price_inr": get_academic_price(cfg.price_inr) if cfg.price_inr > 0 else 0,
            "daily_analyses": cfg.daily_analyses,
            "batch_max_seq": cfg.batch_max_seq,
            "api_calls_per_month": cfg.api_calls_per_month,
            "max_seats": cfg.max_seats,
            "period": cfg.period,
            "description": cfg.description,
        })

    return jsonify({
        "plans": plans,
        "academic_discount_pct": ACADEMIC_DISCOUNT_PCT,
        "currency": "INR",
    }), 200


@payment_bp.route("/api/payments/token-balance", methods=["GET"])
@require_auth
def get_token_balance():
    """Get current user's token balance and subscription status."""
    usage = check_usage(g.user["email"])
    return jsonify(usage), 200


@payment_bp.route("/api/payments/create-order", methods=["POST"])
@require_auth
def create_order():
    """Create a Razorpay order. Server-authoritative pricing — never trusts client amounts."""
    data = request.get_json(silent=True) or {}
    product_id = data.get("product_id") or data.get("plan_id", "")
    raw_quantity = data.get("quantity", 1)

    # Type safety: reject non-string product_id
    if not isinstance(product_id, str):
        return jsonify({"error": "Invalid product_id."}), 400

    # Validate quantity using centralized validator (handles NaN, Inf, float, etc.)
    valid, quantity, err = validate_quantity(raw_quantity)
    if not valid:
        return jsonify({"error": err}), 400

    # Validate
    error = validate_order_request(product_id, quantity)
    if error:
        return jsonify({"error": error}), 400

    # Calculate amount from server-side registry (NEVER from client)
    amount_paise = get_amount_paise(product_id, quantity)

    # Apply academic discount if applicable
    discount = data.get("discount", 0)
    if discount:
        discount = min(int(discount), ACADEMIC_DISCOUNT_PCT)
        amount_paise = int(amount_paise * (100 - discount) / 100)

    # Determine designs to credit for display
    designs = get_designs_for_product(product_id, quantity)

    # Get user info for Razorpay prefill
    user_row = fetch_one(
        "SELECT id, full_name, email FROM users WHERE email = %s",
        (g.user["email"],)
    )
    user_id = user_row["id"]
    import re as _re
    raw_name = user_row.get("full_name") or g.user["email"].split("@")[0]
    user_name = _re.sub(r"[^a-zA-Z\s]", " ", raw_name).strip()
    if len(user_name) < 3:
        user_name = "VigyanLLM User"

    # Create Razorpay order
    receipt = f"pf_{int(time.time())}_{product_id}_{quantity}"
    if not rz_client:
        return jsonify({"error": "Payment service not configured."}), 503
    try:
        rz_order = rz_client.order.create({
            "amount": amount_paise,
            "currency": "INR",
            "receipt": receipt,
            "notes": {
                "email": g.user["email"],
                "user_id": str(user_id),
                "product_id": product_id,
                "quantity": str(quantity),
                "designs": str(designs),
            }
        })
    except Exception as e:
        logger.error("Razorpay order creation failed: %s", e)
        return jsonify({"error": "Payment service unavailable."}), 500

    # Persist order in database
    try:
        execute(
            """INSERT INTO payments
               (user_id, gateway_order_id, amount, currency, status, product_type, tokens_purchased, metadata)
               VALUES (%s, %s, %s, 'INR', 'initiated', %s, %s, %s)""",
            (user_id, rz_order["id"], amount_paise / 100, product_id, designs,
             json.dumps({"quantity": quantity, "receipt": receipt}))
        )
    except Exception:
        logger.warning("Failed to persist payment order (table may not exist)")

    log_action(g.user["email"], "order_created",
               f"Order {rz_order['id']} for {product_id} ({designs} designs), ₹{amount_paise // 100}")

    return jsonify({
        "order_id": rz_order["id"],
        "amount": amount_paise,
        "currency": "INR",
        "key_id": RAZORPAY_KEY_ID,
        "product_id": product_id,
        "tokens": designs,
        "description": f"VigyanLLM: {designs} design(s)",
        "theme": {"color": "#2563EB"},
        "prefill": {
            "name": user_name,
            "email": g.user["email"],
        }
    }), 200


@payment_bp.route("/api/payments/verify-payment", methods=["POST"])
@require_auth
def verify_payment():
    """Verify Razorpay payment signature and credit tokens atomically."""
    data = request.get_json(silent=True) or {}
    razorpay_payment_id = data.get("razorpay_payment_id", "")
    razorpay_order_id = data.get("razorpay_order_id", "")
    razorpay_signature = data.get("razorpay_signature", "")

    if not all([razorpay_payment_id, razorpay_order_id, razorpay_signature]):
        return jsonify({"error": "Missing payment verification fields."}), 400

    # Look up order FIRST to prevent CPU exhaustion on fake payloads (Task 15)
    order = fetch_one(
        """SELECT p.id, p.user_id, p.product_type, p.tokens_purchased, p.status,
                  p.metadata
           FROM payments p
           JOIN users u ON u.id = p.user_id
           WHERE p.gateway_order_id = %s AND u.email = %s""",
        (razorpay_order_id, g.user["email"])
    )

    if not order:
        return jsonify({"error": "Order not found."}), 404

    # Verify HMAC-SHA256 signature (client-side check)
    if not _verify_signature(razorpay_order_id, razorpay_payment_id, razorpay_signature):
        logger.warning("Signature mismatch for order %s by %s", razorpay_order_id, g.user["email"])
        try:
            execute(
                """INSERT INTO system_events (severity, module, message, context)
                   VALUES ('WARNING', 'payments', 'Signature verification failed', %s)""",
                (json.dumps({"order_id": razorpay_order_id, "user_id": g.user.get("user_id")}),)
            )
        except Exception as e:
            logger.debug("Suppressed exception: %s", e)
        return jsonify({"error": "Payment verification failed."}), 400

    # PAY-01 FIX: Server-side confirmation via Razorpay API
    # HMAC alone cannot prove payment was captured — confirm with Razorpay directly
    if not _confirm_payment_server_side(razorpay_payment_id):
        logger.error(
            "Server-side payment confirmation FAILED for order %s payment %s by %s",
            razorpay_order_id, razorpay_payment_id, g.user["email"],
        )
        return jsonify({
            "error": "Payment could not be confirmed with the payment gateway. "
                     "Please wait a few minutes and contact support if your payment was charged.",
            "code": "PAYMENT_UNCONFIRMED",
        }), 400

    # Get quantity from metadata
    metadata = json.loads(order.get("metadata") or "{}") if isinstance(order.get("metadata"), str) else (order.get("metadata") or {})
    quantity = int(metadata.get("quantity", 1))

    # Atomic idempotent credit
    try:
        tokens_credited = _credit_tokens_atomic(
            order["user_id"], razorpay_order_id, order["product_type"],
            quantity, razorpay_payment_id
        )
    except Exception as e:
        logger.error("Token credit failed: %s", e)
        return jsonify({"error": "Token credit failed. Contact support."}), 500

    try:
        log_action(g.user["email"], "payment_verified",
                   f"order={razorpay_order_id} payment={razorpay_payment_id} tokens={tokens_credited}")
    except Exception as e:
        logger.debug("Suppressed exception: %s", e)

    # Return updated balance
    try:
        usage = check_usage(g.user["email"])
    except Exception:
        usage = {"can_run": True, "balance": tokens_credited}
    return jsonify({
        "success": True,
        "tokens_credited": tokens_credited,
        "message": f"{tokens_credited} token(s) credited." if tokens_credited > 0
                   else "Payment already processed.",
        "usage": usage,
    }), 200


@payment_bp.route("/api/payments/webhook", methods=["POST"])
def razorpay_webhook():
    """
    Razorpay server-to-server webhook.
    Always returns 200 to prevent retry storms.
    Uses standalone DB connection (no Flask request auth context).
    """
    raw_body = request.get_data(as_text=True)
    webhook_signature = request.headers.get("X-Razorpay-Signature", "")

    # Verify webhook signature
    expected_sig = hmac.new(
        RAZORPAY_WEBHOOK_SECRET.encode(),
        raw_body.encode(),
        hashlib.sha256
    ).hexdigest()

    # Store webhook regardless of validation
    validation_status = "verified" if hmac.compare_digest(expected_sig, webhook_signature) else "untrusted"

    try:
        event = json.loads(raw_body)
    except Exception:
        return jsonify({"status": "invalid_json"}), 200

    event_type = event.get("event", "")

    # Log webhook to gateway_webhooks table
    try:
        conn = get_db_standalone()
        cur = conn.cursor()
        cur.execute(
            """INSERT INTO gateway_webhooks (raw_payload, event_type, validation_status, http_headers)
               VALUES (%s, %s, %s, %s)""",
            (json.dumps(event), event_type, validation_status,
             json.dumps(dict(request.headers)))
        )
        conn.commit()

        if validation_status == "untrusted":
            logger.warning("Webhook signature failed for event: %s", event_type)
            cur.close()
            put_db_standalone(conn)
            return jsonify({"status": "signature_invalid"}), 200

        # Process payment.captured
        if event_type == "payment.captured":
            payload = event.get("payload", {}).get("payment", {}).get("entity", {})
            order_id = payload.get("order_id", "")
            payment_id = payload.get("id", "")
            notes = payload.get("notes", {})
            email = notes.get("email", "")
            product_id = notes.get("product_id", "")
            quantity = int(notes.get("quantity", "1"))

            if order_id and email:
                # Find user and order
                cur.execute(
                    """SELECT p.id, p.user_id, p.product_type, p.tokens_purchased, p.status
                       FROM payments p
                       JOIN users u ON u.id = p.user_id
                       WHERE p.gateway_order_id = %s AND u.email = %s""",
                    (order_id, email)
                )
                order_row = cur.fetchone()

                if order_row and order_row["status"] in ("initiated", "authorized"):
                    user_id = order_row["user_id"]
                    tokens = order_row["tokens_purchased"] or quantity

                    # Conditional update (idempotent)
                    cur.execute(
                        """UPDATE payments SET status = 'captured', captured_at = NOW(),
                                  gateway_payment_id = %s
                           WHERE gateway_order_id = %s AND status IN ('initiated', 'authorized')""",
                        (payment_id, order_id)
                    )

                    if cur.rowcount > 0:
                        # Credit tokens
                        cur.execute(
                            """UPDATE token_balances
                               SET balance = balance + %s,
                                   total_purchased = total_purchased + %s,
                                   last_credited_at = NOW()
                               WHERE user_id = %s""",
                            (tokens, tokens, user_id)
                        )

                        # Handle subscription
                        if product_id == "base_subscription":
                            cur.execute(
                                """INSERT INTO subscriptions (user_id, is_active, started_at, expires_at, last_renewed_at)
                                   VALUES (%s, TRUE, NOW(), NOW() + INTERVAL '365 days', NOW())
                                   ON CONFLICT (user_id) DO UPDATE SET
                                     is_active = TRUE, expires_at = NOW() + INTERVAL '365 days', last_renewed_at = NOW()""",
                                (user_id,)
                            )

                        conn.commit()
                        logger.info("Webhook: credited %s tokens to user_id=%s for order %s", tokens, user_id, order_id)

        elif event_type == "payment.failed":
            payload = event.get("payload", {}).get("payment", {}).get("entity", {})
            order_id = payload.get("order_id", "")
            if order_id:
                cur.execute(
                    """UPDATE payments SET status = 'failed', failed_at = NOW()
                       WHERE gateway_order_id = %s AND status IN ('initiated', 'authorized')""",
                    (order_id,)
                )
                conn.commit()

        # ── Subscription Events ───────────────────────────────────────────
        elif event_type == "subscription.charged":
            # Monthly renewal — reset quota for next billing cycle
            payload = event.get("payload", {}).get("subscription", {}).get("entity", {})
            sub_id = payload.get("id", "")
            if sub_id:
                cur.execute(
                    """UPDATE subscriptions SET quota_used = 0,
                              dock_quota_used = 0,
                              quota_reset_at = NOW() + INTERVAL '30 days',
                              last_renewed_at = NOW(),
                              expires_at = NOW() + INTERVAL '30 days'
                       WHERE razorpay_subscription_id = %s AND is_active = TRUE""",
                    (sub_id,)
                )
                conn.commit()
                logger.info("Webhook: subscription.charged — quota reset for sub %s", sub_id)

        elif event_type == "subscription.authenticated":
            # New subscription authorized — activate plan
            payload = event.get("payload", {}).get("subscription", {}).get("entity", {})
            sub_id = payload.get("id", "")
            plan_id = payload.get("notes", {}).get("plan_id", "")
            email = payload.get("notes", {}).get("email", "")
            if sub_id and email and plan_id:
                user_row = cur.execute("SELECT id FROM users WHERE email = %s", (email,))
                user_row = cur.fetchone()
                if user_row and plan_id in PRICE_REGISTRY:
                    product = PRICE_REGISTRY[plan_id]
                    expiry_interval = "24 hours" if product.period == "daily" else "30 days"
                    cur.execute(
                        f"""INSERT INTO subscriptions (user_id, is_active, plan_id, plan_type,
                               monthly_quota, dock_monthly_quota, quota_used, dock_quota_used,
                               started_at, expires_at, last_renewed_at, max_seats,
                               razorpay_subscription_id, quota_reset_at)
                           VALUES (%s, TRUE, %s, %s, %s, %s, 0, 0,
                                   NOW(), NOW() + INTERVAL '{expiry_interval}', NOW(), %s, %s, NOW() + INTERVAL '{expiry_interval}')
                           ON CONFLICT (user_id) DO UPDATE SET
                             is_active = TRUE, plan_id = %s, plan_type = %s,
                             monthly_quota = %s, dock_monthly_quota = %s,
                             quota_used = 0, dock_quota_used = 0,
                             expires_at = NOW() + INTERVAL '{expiry_interval}',
                             last_renewed_at = NOW(), max_seats = %s,
                             razorpay_subscription_id = %s,
                             quota_reset_at = NOW() + INTERVAL '{expiry_interval}'""",
                        (user_row["id"], plan_id, plan_id, product.designs_included, product.dock_runs_included,
                         product.max_seats, sub_id,
                         plan_id, plan_id, product.designs_included, product.dock_runs_included,
                         product.max_seats, sub_id)
                    )
                    conn.commit()
                    logger.info("Webhook: subscription.authenticated — %s activated %s", email, plan_id)

        elif event_type in ("subscription.halted", "subscription.cancelled"):
            # Subscription stopped — deactivate
            payload = event.get("payload", {}).get("subscription", {}).get("entity", {})
            sub_id = payload.get("id", "")
            if sub_id:
                cur.execute(
                    """UPDATE subscriptions SET is_active = FALSE
                       WHERE razorpay_subscription_id = %s""",
                    (sub_id,)
                )
                conn.commit()
                logger.info("Webhook: %s — subscription %s deactivated", event_type, sub_id)

        cur.close()
        put_db_standalone(conn)

    except Exception as e:
        logger.error("Webhook processing error: %s", e)
        if "conn" in locals() and conn is not None:
            try:
                put_db_standalone(conn)
            except Exception as e:
                logger.debug("Suppressed exception: %s", e)

    return jsonify({"status": "ok"}), 200


# ══════════════════════════════════════════════════════════════════════════
# ADMIN: Financial Dashboard Endpoints
# ══════════════════════════════════════════════════════════════════════════

@payment_bp.route("/api/payments/financial-summary", methods=["GET"])
@require_admin
def financial_summary():
    """Admin-only: Get P&L, ROI, and user profitability data."""
    pnl = fetch_all("SELECT * FROM v_monthly_pnl LIMIT 12")
    roi = fetch_all("SELECT * FROM v_roi_dashboard LIMIT 12")
    token_econ = fetch_all("SELECT * FROM v_token_economics LIMIT 12")
    admin_costs = fetch_all("SELECT * FROM v_admin_cost_breakdown LIMIT 50")
    top_users = fetch_all("SELECT * FROM v_user_profitability LIMIT 20")

    return jsonify({
        "monthly_pnl": pnl,
        "roi_dashboard": roi,
        "token_economics": token_econ,
        "admin_cost_breakdown": admin_costs,
        "top_users_by_profit": top_users,
    }), 200


@payment_bp.route("/api/payments/revenue-stats", methods=["GET"])
@require_admin
def revenue_stats():
    """
    Admin-only: Real revenue = only money actually paid to us (captured payments).
    Cost = infrastructure cost from pipeline runs (cost_ledger).
    This separates revenue (payments) from cost (pipeline usage).
    """
    # Revenue: only from captured payments (actual money received)
    rev = fetch_one("""
        SELECT
            COALESCE(SUM(amount), 0) AS total_revenue_inr,
            COUNT(*) AS total_payments,
            COALESCE(SUM(tokens_purchased), 0) AS total_tokens_sold
        FROM payments
        WHERE status = 'captured'
    """)

    # Cost: infrastructure cost from all pipeline runs
    cost = fetch_one("""
        SELECT
            COALESCE(SUM(total_cogs_inr), 0) AS total_cost_inr,
            COUNT(*) AS total_operations,
            COALESCE(SUM(CASE WHEN is_billable THEN total_cogs_inr ELSE 0 END), 0) AS paid_user_cost,
            COALESCE(SUM(CASE WHEN NOT is_billable THEN total_cogs_inr ELSE 0 END), 0) AS admin_free_cost
        FROM cost_ledger
    """) or {"total_cost_inr": 0, "total_operations": 0, "paid_user_cost": 0, "admin_free_cost": 0}

    total_rev = float(rev["total_revenue_inr"]) if rev else 0
    total_cost = float(cost["total_cost_inr"]) if cost else 0
    margin = total_rev - total_cost

    return jsonify({
        "revenue": {
            "total_inr": total_rev,
            "payments_count": rev["total_payments"] if rev else 0,
            "tokens_sold": rev["total_tokens_sold"] if rev else 0,
        },
        "cost": {
            "total_inr": total_cost,
            "operations_count": cost["total_operations"],
            "paid_user_cost_inr": float(cost["paid_user_cost"]),
            "admin_free_cost_inr": float(cost["admin_free_cost"]),
        },
        "margin": {
            "gross_profit_inr": margin,
            "margin_percent": round(margin / total_rev * 100, 1) if total_rev > 0 else 0,
        }
    }), 200


@payment_bp.route("/api/payments/user-profitability/<int:user_id>", methods=["GET"])
@require_admin
def user_profitability(user_id: int):
    """Admin-only: Get financial summary for a specific user."""
    summary = fetch_one("SELECT * FROM fn_user_financial_summary(%s)", (user_id,))
    if not summary:
        return jsonify({"error": "User not found."}), 404
    return jsonify(summary), 200


# ══════════════════════════════════════════════════════════════════════════
# PLAN STATUS (for checkout.html frontend)
# ══════════════════════════════════════════════════════════════════════════

@payment_bp.route("/api/payments/status", methods=["GET"])
@payment_bp.route("/api/payment/status", methods=["GET"])
@require_auth
def payment_status():
    """Get current user's plan and usage status (PostgreSQL)."""
    import re
    email = g.user["email"]

    # Get subscription info
    sub = fetch_one(
        """SELECT s.is_active, s.plan_id, s.monthly_quota, s.quota_used,
                  s.quota_reset_at, s.expires_at
           FROM subscriptions s
           JOIN users u ON u.id = s.user_id
           WHERE u.email = %s AND s.is_active = true""",
        (email,)
    )

    # Get academic status
    acad = fetch_one(
        """SELECT status FROM academic_claims WHERE user_id = (SELECT id FROM users WHERE email = %s)""",
        (email,)
    )
    is_academic = bool(acad and acad["status"] == "approved")

    # Auto-detect academic status if not yet set
    if not is_academic:
        try:
            domain = email.split("@")[1].lower()
            _edu_re = r'(^|\.)(edu|ac\.in|edu\.in|ac\.uk|edu\.au|ac\.nz|ac\.jp|edu\.cn|ac\.cn|ac\.kr|edu\.kr|ac\.th|edu\.tw|ac\.za|edu\.mx|ac\.cl|edu\.ar|edu\.sg|edu\.my|edu\.hk|ac\.id|edu\.eg|ac\.ma|edu\.vn|edu\.pk|ac\.ir|edu\.tr|edu\.jo|edu\.lb|ac\.il)(\.[a-z]{2})?$'
            if re.search(_edu_re, domain, re.I) or ".edu." in domain or ".ac." in domain:
                is_academic = True
                try:
                    uid = fetch_one("SELECT id FROM users WHERE email = %s", (email,))
                    if uid:
                        execute("""
                            INSERT INTO academic_claims (user_id, email_edu, status, created_at)
                            VALUES (%s, %s, 'approved', NOW())
                            ON CONFLICT (user_id) DO UPDATE SET status = 'approved'
                        """, (uid["id"], email))
                except Exception:
                    pass
        except Exception:
            pass

    # Get usage info from check_usage
    usage = check_usage(email)

    FREE_DAILY_LIMIT = 5
    plan = "free"
    billing_cycle = "monthly"
    plan_expires_at = 0

    uid_row = fetch_one("SELECT id FROM users WHERE email = %s", (email,))
    user_id = uid_row["id"] if uid_row else None

    if sub:
        plan = sub["plan_id"] or "free"
        if plan:
            plan = plan.split("-")[0] if "-" in plan else plan
        if sub.get("expires_at"):
            try:
                plan_expires_at = int(sub["expires_at"].timestamp())
            except Exception:
                plan_expires_at = 0
        billing_cycle = "monthly"
        if sub["plan_id"] and "-" in (sub["plan_id"] or ""):
            billing_cycle = sub["plan_id"].split("-")[1] if len(sub["plan_id"].split("-")) > 1 else "monthly"

    # Count today's usage from agent_work_logs (works for all tiers)
    today_start = time.time() - (time.time() % 86400)
    today_usage = 0
    if user_id:
        try:
            row = fetch_one(
                """SELECT COUNT(*) AS cnt FROM agent_work_logs
                   WHERE user_id = %s AND completed_at >= to_timestamp(%s)""",
                (user_id, today_start)
            )
            if row:
                today_usage = row["cnt"] or 0
        except Exception:
            pass

    if sub and sub.get("is_active") and plan != "free":
        daily_quota = sub["monthly_quota"] or FREE_DAILY_LIMIT
        daily_used = sub["quota_used"] or 0
        can_analyze = daily_used < daily_quota
        remaining = max(0, daily_quota - daily_used)
    else:
        daily_quota = FREE_DAILY_LIMIT
        daily_used = today_usage
        can_analyze = today_usage < FREE_DAILY_LIMIT
        remaining = max(0, FREE_DAILY_LIMIT - today_usage)

    return jsonify({
        "plan": plan,
        "billing_cycle": billing_cycle,
        "plan_activated_at": 0,
        "plan_expires_at": plan_expires_at,
        "is_academic": is_academic,
        "academic_discount": 30 if is_academic else 0,
        "daily": {
            "can_analyze": can_analyze,
            "used": daily_used,
            "limit": daily_quota,
            "remaining": remaining,
            "tool_used": 0,
            "tool_limit": daily_quota,
        },
        "api": {"can_call": True, "used": 0, "limit": 0},
        "razorpay_key_id": RAZORPAY_KEY_ID,
    }), 200


# ══════════════════════════════════════════════════════════════════════════
# USAGE CHECK & RECORD (for frontend gates)
# ══════════════════════════════════════════════════════════════════════════

@payment_bp.route("/api/usage/check", methods=["GET"])
@require_auth
def usage_check():
    """Check if user can run an analysis (daily limit for free, monthly for paid)."""
    email = g.user["email"]
    today_start = time.time() - (time.time() % 86400)
    uid_row = fetch_one("SELECT id FROM users WHERE email = %s", (email,))
    user_id = uid_row["id"] if uid_row else None

    # Get subscription info
    sub = fetch_one(
        """SELECT s.is_active, s.plan_id, s.monthly_quota, s.quota_used
           FROM subscriptions s JOIN users u ON u.id = s.user_id
           WHERE u.email = %s AND s.is_active = true""",
        (email,)
    )

    FREE_DAILY = 5
    if sub and sub.get("is_active"):
        limit = sub["monthly_quota"] or FREE_DAILY
        used = sub["quota_used"] or 0
        remaining = max(0, limit - used)
        can_analyze = remaining > 0
    else:
        # Free tier: count today's agent_work_logs
        used = 0
        if user_id:
            try:
                row = fetch_one(
                    """SELECT COUNT(*) AS cnt FROM agent_work_logs
                       WHERE user_id = %s AND completed_at >= to_timestamp(%s)""",
                    (user_id, today_start)
                )
                if row:
                    used = row["cnt"] or 0
            except Exception:
                pass
        limit = FREE_DAILY
        remaining = max(0, limit - used)
        can_analyze = used < limit

    return jsonify({
        "can_analyze": can_analyze,
        "used": used,
        "limit": limit,
        "remaining": remaining,
    }), 200


@payment_bp.route("/api/usage/record", methods=["POST"])
@require_auth
def usage_record():
    """Record a completed analysis in agent_work_logs."""
    data = request.get_json(silent=True) or {}
    email = g.user["email"]
    sequences_count = int(data.get("sequences_count", 1))
    uid_row = fetch_one("SELECT id FROM users WHERE email = %s", (email,))
    if not uid_row:
        return jsonify({"error": "User not found"}), 404
    user_id = uid_row["id"]

    try:
        execute(
            """INSERT INTO agent_work_logs (user_id, agent_name, status, completed_at)
               VALUES (%s, 'primer_analysis', 'completed', NOW())""",
            (user_id,)
        )
    except Exception as e:
        logger.error("Failed to record usage: %s", e)
        return jsonify({"error": "Internal server error"}), 500

    return jsonify({"success": True}), 200


# ══════════════════════════════════════════════════════════════════════════
# RESULTS & EXPORT (for dashboard frontend)
# ══════════════════════════════════════════════════════════════════════════

@payment_bp.route("/api/results/save", methods=["POST"])
@require_auth
def save_result():
    """Save a user's analysis result."""
    import json as _json
    data = request.get_json(silent=True) or {}
    email = g.user["email"]
    tool = (data.get("tool") or "").strip()
    if not tool:
        return jsonify({"error": "Tool name required"}), 400
    uid_row = fetch_one("SELECT id FROM users WHERE email = %s", (email,))
    if not uid_row:
        return jsonify({"error": "User not found"}), 404
    user_id = uid_row["id"]
    title = data.get("title") or ""
    job_id = data.get("job_id") or ""
    inputs = _json.dumps(data.get("inputs", {}))
    outputs = _json.dumps(data.get("outputs", {}))
    seq_count = int(data.get("sequences_count", 0))
    try:
        rid = execute(
            """INSERT INTO user_reports (user_id, title, full_result)
               VALUES (%s, %s, %s) RETURNING id""",
            (user_id, title, _json.dumps({"tool": tool, "inputs": inputs, "outputs": outputs, "job_id": job_id, "seq_count": seq_count}))
        )
        if isinstance(rid, (list, tuple)):
            rid = rid[0]
        return jsonify({"success": True, "id": rid}), 200
    except Exception as e:
        logger.error("Failed to save result: %s", e)
        return jsonify({"error": "Internal server error"}), 500


@payment_bp.route("/api/results/list", methods=["GET"])
@require_auth
def list_results():
    """List saved results for the current user."""
    email = g.user["email"]
    limit = min(int(request.args.get("limit", 20)), 100)
    offset = int(request.args.get("offset", 0))
    tool_filter = (request.args.get("tool") or "").strip()
    uid_row = fetch_one("SELECT id FROM users WHERE email = %s", (email,))
    if not uid_row:
        return jsonify({"results": [], "total": 0}), 200
    user_id = uid_row["id"]
    try:
        if tool_filter:
            rows = fetch_all(
                """SELECT id, title, full_result, created_at
                   FROM user_reports WHERE user_id = %s
                   ORDER BY created_at DESC LIMIT %s OFFSET %s""",
                (user_id, limit, offset)
            )
            total = fetch_one(
                "SELECT COUNT(*) AS cnt FROM user_reports WHERE user_id = %s",
                (user_id,)
            )
        else:
            rows = fetch_all(
                """SELECT id, title, full_result, created_at
                   FROM user_reports WHERE user_id = %s
                   ORDER BY created_at DESC LIMIT %s OFFSET %s""",
                (user_id, limit, offset)
            )
            total = fetch_one(
                "SELECT COUNT(*) AS cnt FROM user_reports WHERE user_id = %s",
                (user_id,)
            )
    except Exception as e:
        logger.error("Failed to list results: %s", e)
        return jsonify({"error": "Internal server error"}), 500
    results = []
    for r in rows:
        import json as _json
        try:
            fr = r.get("full_result")
            if isinstance(fr, str):
                full = _json.loads(fr)
            elif isinstance(fr, dict):
                full = fr
            else:
                full = {}
        except Exception:
            full = {}
        results.append({
            "id": r["id"],
            "tool": full.get("tool", ""),
            "title": r.get("title", ""),
            "inputs": full.get("inputs", {}),
            "outputs": full.get("outputs", {}),
            "sequences_count": full.get("seq_count", 0),
            "job_id": full.get("job_id", ""),
            "created_at": str(r["created_at"]) if r.get("created_at") else "",
        })
    return jsonify({"results": results, "total": total["cnt"] if total else 0}), 200


@payment_bp.route("/api/results/delete", methods=["POST"])
@require_auth
def delete_result():
    """Delete a saved result."""
    data = request.get_json(silent=True) or {}
    rid = data.get("id")
    if not rid:
        return jsonify({"error": "Result ID required"}), 400
    email = g.user["email"]
    uid_row = fetch_one("SELECT id FROM users WHERE email = %s", (email,))
    if not uid_row:
        return jsonify({"error": "User not found"}), 404
    user_id = uid_row["id"]
    try:
        row = fetch_one(
            "SELECT id FROM user_reports WHERE id = %s AND user_id = %s",
            (rid, user_id)
        )
        if not row:
            return jsonify({"error": "Result not found"}), 404
        execute("DELETE FROM user_reports WHERE id = %s AND user_id = %s", (rid, user_id))
        return jsonify({"success": True}), 200
    except Exception as e:
        logger.error("Failed to delete result: %s", e)
        return jsonify({"error": "Internal server error"}), 500


@payment_bp.route("/api/export/pdf", methods=["POST"])
@require_auth
def export_pdf():
    """Export analysis results as PDF."""
    import json as _json
    import io
    data = request.get_json(silent=True) or {}
    tool = (data.get("tool") or "analysis").strip()
    inputs = data.get("inputs", {})
    outputs = data.get("outputs", {})
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, f"VigyanLLM - {tool.capitalize()} Report", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.set_font("Helvetica", "", 10)
    pdf.cell(0, 8, f"Generated: {time.strftime('%Y-%m-%d %H:%M')}", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.ln(10)
    if inputs:
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 10, "Input Parameters", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 9)
        for k, v in inputs.items():
            if isinstance(v, (dict, list)):
                v = _json.dumps(v)[:200]
            pdf.multi_cell(0, 5, f"{k}: {v}")
        pdf.ln(5)
    if outputs:
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 10, "Results Summary", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 9)
        if isinstance(outputs, str):
            pdf.multi_cell(0, 5, outputs[:2000])
        elif isinstance(outputs, dict):
            summary = outputs.get("summary", _json.dumps(outputs)[:500])
            if isinstance(summary, str):
                pdf.multi_cell(0, 5, summary[:2000])
            else:
                for k, v in (summary.items() if isinstance(summary, dict) else [(k, v) for k, v in outputs.items()]):
                    if isinstance(v, (dict, list)):
                        v = _json.dumps(v)[:200]
                    pdf.multi_cell(0, 5, f"{k}: {v}")
    return Response(bytes(pdf.output()), mimetype="application/pdf", headers={"Content-Disposition": f"attachment; filename={tool}_report.pdf"}), 200


@payment_bp.route("/api/export/pptx", methods=["POST"])
@require_auth
def export_pptx():
    """Export analysis results as PPTX."""
    import json as _json
    import io
    data = request.get_json(silent=True) or {}
    tool = (data.get("tool") or "analysis").strip()
    outputs = data.get("outputs", {})
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"VigyanLLM - {tool.capitalize()} Report"
    p.font.size = Pt(24)
    if outputs:
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
        tf2 = txBox2.text_frame
        summary = outputs.get("summary", str(outputs)[:500]) if isinstance(outputs, dict) else str(outputs)[:500]
        p2 = tf2.paragraphs[0]
        p2.text = str(summary)[:1000]
        p2.font.size = Pt(14)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Response(buf.read(), mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition": f"attachment; filename={tool}_report.pptx"}), 200
